import io
import logging
from typing import List, Dict, Any
from urllib.parse import urlparse
import requests
import pypdf
from docx import Document
import email
from email import policy
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles processing of different document types"""
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.doc', '.eml', '.msg', '.xlsx', '.csv', '.pptx']


    async def scrape_website(self, url: str, selector: str = None) -> Dict[str, Any]:
        """Scrape a website and return text content (optionally using a CSS selector)"""
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')
            if selector:
                elements = soup.select(selector)
                text = '\n'.join([el.get_text(separator=' ', strip=True) for el in elements])
            else:
                text = soup.get_text(separator=' ', strip=True)
            return {
                'content': text,
                'metadata': {
                    'source': url,
                    'type': 'web',
                    'selector': selector or 'all',
                    'length': len(text)
                }
            }
        except Exception as e:
            logger.error(f"Error scraping website {url}: {str(e)}")
            raise
            
            # Determine file type from URL or content type
            content_type = response.headers.get('content-type', '')
            url_path = urlparse(document_url).path.lower()
            
            # Extract text based on file type
            if 'pdf' in content_type or url_path.endswith('.pdf'):
                return await self._process_pdf(response.content, document_url)
            elif 'word' in content_type or url_path.endswith(('.docx', '.doc')):
                return await self._process_docx(response.content, document_url)
            elif url_path.endswith('.xlsx') or 'excel' in content_type:
                return await self._process_xlsx(response.content, document_url)
            elif url_path.endswith('.csv') or 'csv' in content_type:
                return await self._process_csv(response.content, document_url)
            elif url_path.endswith('.pptx') or 'powerpoint' in content_type:
                return await self._process_pptx(response.content, document_url)
            elif 'email' in content_type or url_path.endswith(('.eml', '.msg')):
                return await self._process_email(response.content, document_url)
            elif 'html' in content_type or url_path.endswith('.html'):
                # General website scraping
                return await self.scrape_website(document_url)
            else:
                # Try to extract as text
                return await self._process_text(response.text, document_url)
    async def _process_xlsx(self, content: bytes, source_url: str) -> Dict[str, Any]:
        """Extract text from Excel XLSX file (all sheets)"""
        try:
            import openpyxl
            from tempfile import NamedTemporaryFile
            with NamedTemporaryFile(delete=True, suffix='.xlsx') as tmp:
                tmp.write(content)
                tmp.flush()
                wb = openpyxl.load_workbook(tmp.name, data_only=True)
                text_content = ""
                sheets_info = []
                for sheet in wb.worksheets:
                    sheet_text = []
                    for row in sheet.iter_rows(values_only=True):
                        row_text = ' '.join([str(cell) if cell is not None else '' for cell in row])
                        sheet_text.append(row_text)
                    sheet_text_str = '\n'.join(sheet_text)
                    text_content += f"\n\n--- Sheet: {sheet.title} ---\n{sheet_text_str}"
                    sheets_info.append({'sheet': sheet.title, 'rows': len(sheet_text)})
            return {
                'content': text_content,
                'metadata': {
                    'source': source_url,
                    'type': 'xlsx',
                    'sheets': [s['sheet'] for s in sheets_info],
                    'sheets_info': sheets_info
                }
            }
        except Exception as e:
            logger.error(f"Error processing XLSX: {str(e)}")
            raise

    async def _process_csv(self, content: bytes, source_url: str) -> Dict[str, Any]:
        """Extract text from CSV file"""
        try:
            import csv
            import io as _io
            text_content = ""
            reader = csv.reader(_io.StringIO(content.decode('utf-8')))
            rows = list(reader)
            for row in rows:
                text_content += ', '.join(row) + '\n'
            return {
                'content': text_content,
                'metadata': {
                    'source': source_url,
                    'type': 'csv',
                    'rows': len(rows)
                }
            }
        except Exception as e:
            logger.error(f"Error processing CSV: {str(e)}")
            raise

    async def _process_pptx(self, content: bytes, source_url: str) -> Dict[str, Any]:
        """Extract text from PowerPoint PPTX file (all slides)"""
        try:
            from pptx import Presentation
            from tempfile import NamedTemporaryFile
            with NamedTemporaryFile(delete=True, suffix='.pptx') as tmp:
                tmp.write(content)
                tmp.flush()
                prs = Presentation(tmp.name)
                text_content = ""
                slides_info = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    slide_text_str = '\n'.join(slide_text)
                    text_content += f"\n\n--- Slide {i+1} ---\n{slide_text_str}"
                    slides_info.append({'slide': i+1, 'length': len(slide_text_str)})
            return {
                'content': text_content,
                'metadata': {
                    'source': source_url,
                    'type': 'pptx',
                    'slides': len(slides_info),
                    'slides_info': slides_info
                }
            }
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise
    
    async def _process_pdf(self, content: bytes, source_url: str) -> Dict[str, Any]:
        """Extract text from PDF"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = pypdf.PdfReader(pdf_file)
            
            text_content = ""
            pages_info = []
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                text_content += f"\n\n--- Page {page_num} ---\n{page_text}"
                pages_info.append({
                    'page_number': page_num,
                    'text_length': len(page_text),
                    'text': page_text
                })
            
            return {
                'content': text_content,
                'metadata': {
                    'source': source_url,
                    'type': 'pdf',
                    'total_pages': len(pdf_reader.pages),
                    'pages_info': pages_info
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            raise
    
    async def _process_docx(self, content: bytes, source_url: str) -> Dict[str, Any]:
        """Extract text from DOCX"""
        try:
            doc_file = io.BytesIO(content)
            document = Document(doc_file)
            
            text_content = ""
            paragraphs_info = []
            
            for i, paragraph in enumerate(document.paragraphs):
                text_content += f"{paragraph.text}\n"
                paragraphs_info.append({
                    'paragraph_number': i + 1,
                    'text': paragraph.text
                })
            
            return {
                'content': text_content,
                'metadata': {
                    'source': source_url,
                    'type': 'docx',
                    'total_paragraphs': len(document.paragraphs),
                    'paragraphs_info': paragraphs_info
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing DOCX: {str(e)}")
            raise
    
    async def _process_email(self, content: bytes, source_url: str) -> Dict[str, Any]:
        """Extract text from email"""
        try:
            # Parse email using built-in email library
            msg = email.message_from_bytes(content, policy=policy.default)
            
            text_content = f"Subject: {msg.get('Subject', 'No Subject')}\n"
            text_content += f"From: {msg.get('From', 'Unknown')}\n"
            text_content += f"To: {msg.get('To', 'Unknown')}\n"
            text_content += f"Date: {msg.get('Date', 'Unknown')}\n\n"
            
            # Extract body content
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text_content += f"Body:\n{part.get_content()}\n"
                        break
                    elif part.get_content_type() == "text/html":
                        # Extract text from HTML
                        html_content = part.get_content()
                        soup = BeautifulSoup(html_content, 'html.parser')
                        text_content += f"Body:\n{soup.get_text()}\n"
                        break
            else:
                if msg.get_content_type() == "text/plain":
                    text_content += f"Body:\n{msg.get_content()}\n"
                elif msg.get_content_type() == "text/html":
                    # Extract text from HTML
                    html_content = msg.get_content()
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text_content += f"Body:\n{soup.get_text()}\n"
            
            return {
                'content': text_content,
                'metadata': {
                    'source': source_url,
                    'type': 'email',
                    'subject': msg.get('Subject', 'No Subject'),
                    'from': msg.get('From', 'Unknown'),
                    'to': msg.get('To', 'Unknown'),
                    'date': msg.get('Date', 'Unknown')
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing email: {str(e)}")
            raise
    
    async def _process_text(self, content: str, source_url: str) -> Dict[str, Any]:
        """Process plain text content"""
        return {
            'content': content,
            'metadata': {
                'source': source_url,
                'type': 'text',
                'length': len(content)
            }
        }
